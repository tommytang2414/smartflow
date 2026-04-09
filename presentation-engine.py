"""McKinsey-style Presentation Engine for SmartFlow."""

import json
import sys
import argparse
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import io

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x05, 0x1C, 0x2A)
NAVY2       = RGBColor(0x1B, 0x2A, 0x4A)
TEAL        = RGBColor(0x30, 0xA3, 0xDA)
TEAL2       = RGBColor(0x4E, 0xCD, 0xC4)
TEXT_DARK   = RGBColor(0x2D, 0x34, 0x36)
TEXT_LIGHT  = RGBColor(0x63, 0x6E, 0x72)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG     = RGBColor(0xF0, 0xF2, 0xF5)
GREEN       = RGBColor(0x2E, 0x8B, 0x57)
AMBER       = RGBColor(0xDA, 0x9B, 0x30)
RED         = RGBColor(0xDC, 0x14, 0x3C)
LIGHT_GREY  = RGBColor(0xE8, 0xEC, 0xF0)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def rgb(color): return color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_name="Calibri", font_size=12, bold=False, italic=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True, v_anchor=None):
    from pptx.util import Pt
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                       font_name="Calibri", font_size=11, bold=False,
                       color=TEXT_DARK, line_spacing=None):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


# ── Slide builders ────────────────────────────────────────────────────────────

def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.title  # may not exist on blank

    # Full navy background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)

    # Teal accent bar
    add_rect(slide, 0, Inches(4.5), SLIDE_W, Inches(0.06), fill_color=TEAL)

    # Title
    add_text(slide, data["title"],
             Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.4),
             font_name="Montserrat", font_size=36, bold=True, color=WHITE,
             align=PP_ALIGN.LEFT)

    # Subtitle
    add_text(slide, data.get("subtitle", ""),
             Inches(0.8), Inches(3.3), Inches(9), Inches(0.7),
             font_size=16, color=TEAL, align=PP_ALIGN.LEFT)

    # Date / meta
    meta = data.get("meta", datetime.today().strftime("%B %Y"))
    add_text(slide, meta,
             Inches(0.8), Inches(6.6), Inches(6), Inches(0.5),
             font_size=11, color=TEXT_LIGHT, align=PP_ALIGN.LEFT)


def build_section_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Left navy panel
    add_rect(slide, 0, 0, Inches(4.5), SLIDE_H, fill_color=NAVY2)
    # Right white
    add_rect(slide, Inches(4.5), 0, Inches(8.83), SLIDE_H, fill_color=WHITE)
    # Teal vertical accent
    add_rect(slide, Inches(4.5), 0, Inches(0.07), SLIDE_H, fill_color=TEAL)

    # Section number
    add_text(slide, data.get("section_number", ""),
             Inches(0.5), Inches(2.5), Inches(3.5), Inches(1.2),
             font_name="Montserrat", font_size=60, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER)

    # Title on right
    add_text(slide, data["title"],
             Inches(5.0), Inches(2.8), Inches(7.8), Inches(1.5),
             font_name="Montserrat", font_size=28, bold=True, color=NAVY,
             align=PP_ALIGN.LEFT)

    if data.get("subtitle"):
        add_text(slide, data["subtitle"],
                 Inches(5.0), Inches(4.1), Inches(7.8), Inches(0.7),
                 font_size=14, color=TEXT_LIGHT)


def build_kpi_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)

    # Top navy bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill_color=NAVY)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), fill_color=TEAL)

    add_text(slide, data["title"],
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
             font_name="Montserrat", font_size=20, bold=True, color=WHITE)

    kpis = data.get("kpis", [])
    n = len(kpis)
    card_w = Inches(12.3 / n)
    margin_left = Inches(0.5)
    card_top = Inches(1.5)
    card_h = Inches(4.8)
    gap = Inches(0.2)

    for i, kpi in enumerate(kpis):
        cx = margin_left + i * (card_w + gap)
        add_rect(slide, cx, card_top, card_w, card_h, fill_color=CARD_BG)

        # Colour top accent
        status = kpi.get("status", "")
        accent = TEAL if status == "green" else AMBER if status == "amber" else RED if status == "red" else TEAL
        add_rect(slide, cx, card_top, card_w, Inches(0.07), fill_color=accent)

        # Value
        add_text(slide, kpi.get("value", ""),
                 cx + Inches(0.15), card_top + Inches(0.4), card_w - Inches(0.3), Inches(1.6),
                 font_name="Montserrat", font_size=40, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER)

        # Label
        add_text(slide, kpi.get("label", ""),
                 cx + Inches(0.1), card_top + Inches(1.9), card_w - Inches(0.2), Inches(0.7),
                 font_size=12, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)

        # Sub
        if kpi.get("sub"):
            add_text(slide, kpi["sub"],
                     cx + Inches(0.1), card_top + Inches(2.55), card_w - Inches(0.2), Inches(1.8),
                     font_size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, wrap=True)


def build_text_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill_color=NAVY)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), fill_color=TEAL)

    add_text(slide, data["title"],
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
             font_name="Montserrat", font_size=20, bold=True, color=WHITE)

    bullets = data.get("bullets", [])
    top = Inches(1.5)
    for b in bullets:
        # Bullet symbol
        add_text(slide, "▪",
                 Inches(0.5), top, Inches(0.3), Inches(0.5),
                 font_size=12, color=TEAL, bold=True)
        add_text(slide, b,
                 Inches(0.85), top, Inches(11.8), Inches(0.55),
                 font_size=12, color=TEXT_DARK)
        top += Inches(0.65)

    if data.get("note"):
        add_rect(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.65), fill_color=CARD_BG)
        add_text(slide, f"Note: {data['note']}",
                 Inches(0.65), Inches(6.55), Inches(12), Inches(0.55),
                 font_size=9, color=TEXT_LIGHT, italic=True)


def build_two_column_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill_color=NAVY)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), fill_color=TEAL)

    add_text(slide, data["title"],
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
             font_name="Montserrat", font_size=20, bold=True, color=WHITE)

    left_col = data.get("left", {})
    right_col = data.get("right", {})

    # Left column
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(6.1), Inches(5.8), fill_color=CARD_BG)
    add_text(slide, left_col.get("header", ""),
             Inches(0.55), Inches(1.45), Inches(5.8), Inches(0.55),
             font_size=13, bold=True, color=TEAL)
    top = Inches(2.1)
    for b in left_col.get("bullets", []):
        add_text(slide, "▪",
                 Inches(0.6), top, Inches(0.25), Inches(0.5), font_size=11, color=TEAL)
        add_text(slide, b,
                 Inches(0.9), top, Inches(5.4), Inches(0.55), font_size=11, color=TEXT_DARK)
        top += Inches(0.6)

    # Right column
    add_rect(slide, Inches(6.83), Inches(1.35), Inches(6.1), Inches(5.8), fill_color=CARD_BG)
    add_text(slide, right_col.get("header", ""),
             Inches(7.0), Inches(1.45), Inches(5.8), Inches(0.55),
             font_size=13, bold=True, color=TEAL)
    top = Inches(2.1)
    for b in right_col.get("bullets", []):
        add_text(slide, "▪",
                 Inches(7.0), top, Inches(0.25), Inches(0.5), font_size=11, color=TEAL)
        add_text(slide, b,
                 Inches(7.3), top, Inches(5.4), Inches(0.55), font_size=11, color=TEXT_DARK)
        top += Inches(0.6)

    # Centre divider
    add_rect(slide, Inches(6.6), Inches(1.35), Inches(0.06), Inches(5.8), fill_color=TEAL)


def build_table_slide(prs, data):
    from pptx.util import Pt
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill_color=NAVY)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), fill_color=TEAL)

    add_text(slide, data["title"],
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
             font_name="Montserrat", font_size=20, bold=True, color=WHITE)

    headers = data.get("headers", [])
    rows = data.get("rows", [])
    if not headers or not rows:
        return

    n_cols = len(headers)
    n_rows = len(rows)
    tbl_left = Inches(0.5)
    tbl_top = Inches(1.35)
    tbl_w = Inches(12.33)
    tbl_h = Inches(5.8)
    col_w = tbl_w / n_cols
    row_h_header = Inches(0.55)
    row_h = (tbl_h - row_h_header) / n_rows

    # Header row
    for j, h in enumerate(headers):
        add_rect(slide, tbl_left + j * col_w, tbl_top, col_w, row_h_header, fill_color=NAVY2)
        add_text(slide, h,
                 tbl_left + j * col_w + Inches(0.08), tbl_top + Inches(0.08),
                 col_w - Inches(0.16), row_h_header - Inches(0.1),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        bg = CARD_BG if i % 2 == 0 else WHITE
        for j, cell in enumerate(row):
            add_rect(slide, tbl_left + j * col_w, tbl_top + row_h_header + i * row_h,
                     col_w, row_h, fill_color=bg)
            # Status dot for last column
            if j == n_cols - 1 and cell in ("✓ Done", "⚡ In Progress", "○ Planned"):
                dot_color = GREEN if "Done" in cell else AMBER if "Progress" in cell else TEXT_LIGHT
                add_text(slide, cell,
                         tbl_left + j * col_w + Inches(0.08),
                         tbl_top + row_h_header + i * row_h + Inches(0.05),
                         col_w - Inches(0.16), row_h - Inches(0.1),
                         font_size=10, color=dot_color, bold=True, align=PP_ALIGN.CENTER)
            else:
                add_text(slide, str(cell),
                         tbl_left + j * col_w + Inches(0.08),
                         tbl_top + row_h_header + i * row_h + Inches(0.05),
                         col_w - Inches(0.16), row_h - Inches(0.1),
                         font_size=10, color=TEXT_DARK, align=PP_ALIGN.CENTER)


def build_chart_slide(prs, data):
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill_color=NAVY)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), fill_color=TEAL)

    add_text(slide, data["title"],
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
             font_name="Montserrat", font_size=20, bold=True, color=WHITE)

    chart_data = ChartData()
    chart_data.categories = data.get("categories", [])
    for series in data.get("series", []):
        chart_data.add_series(series["name"], series["values"])

    chart_type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    ct = chart_type_map.get(data.get("chart_type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart = slide.shapes.add_chart(
        ct, Inches(0.5), Inches(1.35), Inches(12.33), Inches(5.8), chart_data
    ).chart

    chart.has_legend = len(data.get("series", [])) > 1
    chart.has_title = False

    # Style series colours
    colors = [TEAL, NAVY2, AMBER, GREEN, RED]
    for i, series in enumerate(chart.series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = colors[i % len(colors)]


def build_timeline_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill_color=NAVY)
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.05), fill_color=TEAL)

    add_text(slide, data["title"],
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.7),
             font_name="Montserrat", font_size=20, bold=True, color=WHITE)

    items = data.get("items", [])
    n = len(items)
    if not n:
        return

    # Horizontal timeline line
    line_y = Inches(3.5)
    add_rect(slide, Inches(0.6), line_y, Inches(12.1), Inches(0.04), fill_color=LIGHT_GREY)

    step = Inches(12.1) / n
    for i, item in enumerate(items):
        cx = Inches(0.6) + i * step + step / 2
        status = item.get("status", "planned")
        dot_color = GREEN if status == "done" else TEAL if status == "active" else LIGHT_GREY
        dot_r = Inches(0.18)
        add_rect(slide, cx - dot_r, line_y - dot_r + Inches(0.02), dot_r * 2, dot_r * 2, fill_color=dot_color)

        # Phase label above
        add_text(slide, item.get("phase", ""),
                 cx - Inches(0.9), Inches(1.6), Inches(1.8), Inches(0.45),
                 font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        # Title above
        add_text(slide, item.get("title", ""),
                 cx - Inches(0.9), Inches(2.0), Inches(1.8), Inches(1.2),
                 font_size=10, color=TEXT_DARK, align=PP_ALIGN.CENTER, wrap=True)

        # Detail below
        add_text(slide, item.get("detail", ""),
                 cx - Inches(0.9), Inches(3.75), Inches(1.8), Inches(2.5),
                 font_size=9, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, wrap=True)

        # Status badge
        badge_c = GREEN if status == "done" else TEAL if status == "active" else TEXT_LIGHT
        badge_t = "DONE" if status == "done" else "ACTIVE" if status == "active" else "PLANNED"
        add_text(slide, badge_t,
                 cx - Inches(0.4), Inches(6.2), Inches(0.8), Inches(0.3),
                 font_size=8, bold=True, color=badge_c, align=PP_ALIGN.CENTER)


def build_closing_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)
    add_rect(slide, 0, Inches(3.6), SLIDE_W, Inches(0.06), fill_color=TEAL)

    add_text(slide, data.get("title", "Next Steps"),
             Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             font_name="Montserrat", font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    items = data.get("items", [])
    top = Inches(3.85)
    for item in items:
        add_text(slide, "→",
                 Inches(2.5), top, Inches(0.4), Inches(0.5),
                 font_size=14, color=TEAL, bold=True)
        add_text(slide, item,
                 Inches(3.0), top, Inches(7.5), Inches(0.5),
                 font_size=13, color=WHITE)
        top += Inches(0.6)

    if data.get("footer"):
        add_text(slide, data["footer"],
                 Inches(1), Inches(6.7), Inches(11), Inches(0.4),
                 font_size=10, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

BUILDERS = {
    "title": build_title_slide,
    "section": build_section_slide,
    "kpi": build_kpi_slide,
    "text": build_text_slide,
    "two_column": build_two_column_slide,
    "table": build_table_slide,
    "chart": build_chart_slide,
    "timeline": build_timeline_slide,
    "closing": build_closing_slide,
}


def build(json_path, output_path, theme="mckinsey"):
    with open(json_path, "r", encoding="utf-8") as f:
        deck = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for slide_data in deck["slides"]:
        slide_type = slide_data.get("type", "text")
        builder = BUILDERS.get(slide_type)
        if builder:
            builder(prs, slide_data)
        else:
            print(f"Unknown slide type: {slide_type}, skipping")

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("json_file")
    parser.add_argument("--output", "-o", default="output.pptx")
    parser.add_argument("--theme", default="mckinsey")
    args = parser.parse_args()
    build(args.json_file, args.output, args.theme)
